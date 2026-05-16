#!/usr/bin/env python3
"""Generate LLM comparison PPTX - Light business style."""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

OUTPUT = os.path.expanduser("~/Desktop/llm-compare-ppt/全球大模型全景对比.pptx")

with open(os.path.expanduser("~/Desktop/llm-compare-ppt/research-data.json"), "r") as f:
    data = json.load(f)
benchmarks = data["benchmarks"]

# ── Color palette – light business ──
BG_WHITE = RGBColor(255, 255, 255)
BG_LIGHT = RGBColor(245, 246, 250)
BG_CARD = RGBColor(255, 255, 255)
TEXT_MAIN = RGBColor(30, 32, 38)
TEXT_SEC = RGBColor(100, 110, 130)
TEXT_LIGHT = RGBColor(150, 160, 180)
ACCENT_BLUE = RGBColor(25, 118, 210)
ACCENT_GREEN = RGBColor(46, 125, 50)
ACCENT_ORANGE = RGBColor(237, 108, 2)
ACCENT_PURPLE = RGBColor(106, 27, 154)
BORDER_LIGHT = RGBColor(220, 225, 235)
HEADER_BG = RGBColor(25, 118, 210)

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Microsoft YaHei"
    p.alignment = align
    return tf

def multi_tb(slide, left, top, width, height, items, default_size=12, default_color=TEXT_MAIN):
    """Add textbox with multiple paragraphs."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, sz, clr, bld = item, default_size, default_color, False
        else:
            text = item.get("text",""); sz = item.get("size", default_size)
            clr = item.get("color", default_color); bld = item.get("bold", False)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz); p.font.bold = bld; p.font.color.rgb = clr; p.font.name = "Microsoft YaHei"
    return tf

def card(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    """Card shape with optional shadow-like border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color or BORDER_LIGHT
    shape.line.width = Pt(1) if border_color else Pt(0.5)
    shape.shadow.inherit = False
    return shape

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ==== SLIDE 1: COVER ── pure white, blue accent stripe ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)
# Top accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()
# Title
tb(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
   "全球大模型全景对比", 40, True, TEXT_MAIN, PP_ALIGN.CENTER)
tb(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
   "Global LLM Landscape · 2025-2026", 18, False, TEXT_SEC, PP_ALIGN.CENTER)
# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_BLUE; div.line.fill.background()
tb(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
   "中美大模型 · 能力 · 价格 · 多模态 · 市场全面对比", 14, False, TEXT_SEC, PP_ALIGN.CENTER)
tb(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
   "2026年4月 · 数据来源：各模型官方技术报告及第三方评测", 11, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# ==== SLIDE 2: KPI DASHBOARD ── light bg, 4 white cards ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
tb(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.6), "行业关键指标", 26, True, TEXT_MAIN)

top_elo = max(benchmarks, key=lambda x:x["elo"])
top_mmlu = max(benchmarks, key=lambda x:x["mmlu"])
top_code = max(benchmarks, key=lambda x:x["humaneval"])
top_math = max(benchmarks, key=lambda x:x["math"])

kpis = [
    (str(top_elo["elo"]), "Chatbot Arena Elo", f"{top_elo['model']} · {top_elo['company']}", ACCENT_BLUE),
    (f"{top_mmlu['mmlu']}%", "MMLU 知识推理", f"{top_mmlu['model']} · {top_mmlu['company']}", ACCENT_GREEN),
    (f"{top_code['humaneval']}%", "HumanEval 编程", f"{top_code['model']} · {top_code['company']}", ACCENT_PURPLE),
    (f"{top_math['math']}%", "MATH 数学推理", f"{top_math['model']} · {top_math['company']}", ACCENT_ORANGE),
]
positions = [(0.5, 1.2), (6.8, 1.2), (0.5, 4.0), (6.8, 4.0)]

for i, (val, label, sub, accent) in enumerate(kpis):
    l, t = positions[i]
    card(slide, Inches(l), Inches(t), Inches(5.8), Inches(2.3), BG_CARD, BORDER_LIGHT)
    # Value
    tb(slide, Inches(l+0.3), Inches(t+0.3), Inches(3), Inches(1.0),
       val, 36, True, accent)
    # Label
    tb(slide, Inches(l+0.3), Inches(t+1.3), Inches(5), Inches(0.4),
       label, 14, True, TEXT_MAIN)
    # Sub-label
    tb(slide, Inches(l+0.3), Inches(t+1.7), Inches(5), Inches(0.4),
       sub, 11, False, TEXT_SEC)

# ==== SLIDE 3: FULL TABLE ── white bg, clean table ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)
tb(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6), "全模型 Benchmark 对比", 24, True, TEXT_MAIN)

rows = len(benchmarks) + 1
tbl = slide.shapes.add_table(rows, 7, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8)).table

headers = ["模型", "MMLU", "HumanEval", "MATH", "GPQA", "Elo", "上下文"]
for j, h in enumerate(headers):
    c = tbl.cell(0, j); c.text = h
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255); p.font.name = "Microsoft YaHei"
    c.fill.solid(); c.fill.fore_color.rgb = HEADER_BG

for i, m in enumerate(benchmarks):
    flag = "🔵 " if "美国" in m["country"] else "🟢 "
    vals = [flag + m["model"], f"{m['mmlu']}%", f"{m['humaneval']}%", f"{m['math']}%", f"{m['gpqa']}%", str(m["elo"]), m["context"]]
    for j, v in enumerate(vals):
        c = tbl.cell(i+1, j); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10); p.font.name = "Microsoft YaHei"; p.font.color.rgb = TEXT_MAIN
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(240, 245, 250) if i % 2 == 1 else RGBColor(255, 255, 255)

tb(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
   "🔵 美国模型 🟢 中国模型 · 共12款旗舰模型完整对比", 10, False, TEXT_SEC)

# ==== SLIDE 4: US vs CN ── light cards ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
tb(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6), "🇺🇸 vs 🇨🇳 · 格局对比", 24, True, TEXT_MAIN)

# US card
card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), RGBColor(227, 242, 253), RGBColor(187, 222, 251))
multi_tb(slide, Inches(0.8), Inches(1.3), Inches(5.2), Inches(5.2), [
    {"text": "🇺🇸 美国阵营", "size": 22, "bold": True, "color": RGBColor(21, 101, 192)},
    {"text": "", "size": 8},
    {"text": "OpenAI: ChatGPT 600M+ MAU · $315B估值", "size": 13},
    {"text": "Anthropic: 编程能力最强 · 企业安全标杆", "size": 13},
    {"text": "Google: 多模态最强 · 1M上下文领先", "size": 13},
    {"text": "定价高 ($2.5~$15/M input)", "size": 13},
    {"text": "全球市场份额 ~75%", "size": 13},
    {"text": "", "size": 8},
    {"text": "优势：生态 · 品牌 · 融资 · 多模态", "size": 14, "bold": True, "color": RGBColor(21, 101, 192)},
])

# CN card
card(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5), RGBColor(232, 245, 233), RGBColor(200, 230, 201))
multi_tb(slide, Inches(7.1), Inches(1.3), Inches(5.2), Inches(5.2), [
    {"text": "🇨🇳 中国阵营", "size": 22, "bold": True, "color": RGBColor(46, 125, 50)},
    {"text": "", "size": 8},
    {"text": "DeepSeek V4: Elo 1395 多项登顶", "size": 13},
    {"text": "训练仅 $5.6M (比美国低20-40x)", "size": 13},
    {"text": "定价仅为美国 1/5~1/20", "size": 13},
    {"text": "开源生态强大 (Qwen/DeepSeek)", "size": 13},
    {"text": "全球份额 ~20-25%", "size": 13},
    {"text": "", "size": 8},
    {"text": "优势：效率 · 价格 · 开源 · 推理", "size": 14, "bold": True, "color": RGBColor(46, 125, 50)},
])

# ==== SLIDE 5: TRAINING + VALUATION ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)
tb(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6), "训练成本 & 公司估值", 24, True, TEXT_MAIN)

# Training cost table
t2 = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.2), Inches(5.5), Inches(4.5)).table
cost_rows = [
    ["模型", "成本", "效率"],
    ["GPT-4", "$100-200M", "基准线"],
    ["Claude 4 Opus", "~$100M", "~同等"],
    ["DeepSeek V3", "$5.6M", "20-40x 更低"],
    ["DeepSeek V4", "~$15M", "5-10x 更低"],
    ["Llama 3.1 405B", "$100M+", "~同等"],
]
for i, rd in enumerate(cost_rows):
    for j, v in enumerate(rd):
        c = t2.cell(i, j); c.text = v
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(11); p.font.name = "Microsoft YaHei"
        if i == 0:
            p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
            c.fill.solid(); c.fill.fore_color.rgb = HEADER_BG
        else:
            p.font.color.rgb = TEXT_MAIN
            c.fill.solid()
            if "DeepSeek" in rd[0]:
                c.fill.fore_color.rgb = RGBColor(232, 245, 233)
                p.font.color.rgb = ACCENT_GREEN
                if j > 0: p.font.bold = True
            else:
                c.fill.fore_color.rgb = RGBColor(255,255,255)

# Valuation table
t3 = slide.shapes.add_table(8, 2, Inches(6.8), Inches(1.2), Inches(5.8), Inches(4.5)).table
val_rows = [
    ["公司", "估值"],
    ["🇺🇸 OpenAI", "$315B"],
    ["🇺🇸 Anthropic", "$61.5B"],
    ["🇺🇸 xAI", "$75B"],
    ["🇨🇳 DeepSeek", "~$5-10B"],
    ["🇨🇳 智谱AI", "~$20B"],
    ["🇨🇳 月之暗面", "~$20B"],
    ["🇨🇳 MiniMax", "~$15B"],
]
for i, rd in enumerate(val_rows):
    for j, v in enumerate(rd):
        c = t3.cell(i, j); c.text = v
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(11); p.font.name = "Microsoft YaHei"
        if i == 0:
            p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
            c.fill.solid(); c.fill.fore_color.rgb = HEADER_BG
        else:
            p.font.color.rgb = TEXT_MAIN
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(245, 246, 250) if i % 2 == 0 else RGBColor(255,255,255)
            if "🇺🇸" in rd[0] and j == 0:
                p.font.color.rgb = ACCENT_BLUE
            if j == 1 and "🇺🇸" in rd[0]:
                p.font.bold = True
                p.font.color.rgb = ACCENT_BLUE

tb(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
   "DeepSeek 自筹$100M+ 挑战 OpenAI $30B+ 融资 · 效率是最大竞争壁垒", 10, False, TEXT_SEC)

# ==== SLIDE 6: CONCLUSION ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
tb(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "关键结论", 28, True, TEXT_MAIN, PP_ALIGN.CENTER)

conclusions = [
    ("1", "中美差距快速缩小 — Benchmark上中国旗舰已追平甚至领先", ACCENT_BLUE),
    ("2", "成本效率是中国核心竞争力 — 训练成本低5-40倍", ACCENT_GREEN),
    ("3", "价格战由中国主导 — API价格为美国1/5~1/20，加速全球普及", ACCENT_ORANGE),
    ("4", "估值鸿沟巨大 — OpenAI $315B vs DeepSeek ~$10B", RGBColor(200, 35, 35)),
    ("5", "开源重塑格局 — Qwen和DeepSeek开源模型推动全球AI创新", ACCENT_PURPLE),
]
for i, (num, text, accent) in enumerate(conclusions):
    y = Inches(1.8 + i * 1.05)
    # Card
    card(slide, Inches(1.2), y, Inches(10.8), Inches(0.8), BG_CARD, BORDER_LIGHT)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), y+Inches(0.15), Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = accent; badge.line.fill.background()
    tf_b = badge.text_frame; tf_b.paragraphs[0].text = num
    tf_b.paragraphs[0].font.size = Pt(16); tf_b.paragraphs[0].font.bold = True
    tf_b.paragraphs[0].font.color.rgb = RGBColor(255,255,255); tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Text
    tb(slide, Inches(2.2), y+Inches(0.15), Inches(9.5), Inches(0.5), text, 15, False, TEXT_MAIN)

prs.save(OUTPUT)
sz = os.path.getsize(OUTPUT)
print(f"✅ {OUTPUT}", flush=True)
print(f"   {sz} bytes, {len(prs.slides)} slides", flush=True)
