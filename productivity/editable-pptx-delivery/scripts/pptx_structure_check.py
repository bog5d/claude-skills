#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PPTX 交付结构自检器（交付前必跑）。

用途：堵死「整页贴图交差」这条后门。任何 .pptx 交付给波总/投资人前跑一次，
逐页统计真文本框/图片对象/全页图，输出结构化结论。

判定规则：
  - 全页图片 = PICTURE(shape_type==13) 宽高同时超过页面 85%。
  - 若某页存在全页图且该页无文本框 → 该页判「整页贴图」，FAIL。
  - 验收优先级（波总定）：可编辑性 > 内容逻辑 > 图片真实性 > 视觉完成度。

用法:
  python3 pptx_structure_check.py <input.pptx> [--page N] [--json]
    --page N   只检查第 N 页 (1-based)
    --json     输出 JSON（便于脚本间传递）

退出码: 0 = 通过（无整页贴图页）; 1 = 存在整页贴图页。
"""
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Emu


def check(path, page_only=None):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    report = {"slide_count": len(prs.slides), "slides": [], "failed": False}

    for idx, slide in enumerate(prs.slides, 1):
        if page_only and idx != page_only:
            continue
        pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        txts = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        fullpage = []
        maxpic = None
        for s in pics:
            w, h = Emu(s.width).inches, Emu(s.height).inches
            if maxpic is None or (w * h) > (maxpic[0] * maxpic[1]):
                maxpic = (w, h)
            if s.width > sw * 0.85 and s.height > sh * 0.85:
                fullpage.append(s.name)
        is_fullpage = bool(fullpage and not txts)
        if is_fullpage:
            report["failed"] = True
        report["slides"].append({
            "page": idx,
            "text_boxes": len(txts),
            "pictures": len(pics),
            "shapes": len(slide.shapes),
            "full_page_image": is_fullpage,
            "full_page_names": fullpage,
            "max_pic_inches": list(maxpic) if maxpic else None,
        })
    return report


def main():
    ap = argparse.ArgumentParser(description="PPTX 交付结构自检")
    ap.add_argument("input", help=".pptx 路径")
    ap.add_argument("--page", type=int, default=None, help="只检查某页 (1-based)")
    ap.add_argument("--json", action="store_true", help="输出 JSON")
    args = ap.parse_args()

    rep = check(args.input, args.page)
    if args.json:
        print(json.dumps(rep, ensure_ascii=False, indent=2))
    else:
        print("=" * 60)
        print(f"PPTX 结构自检: {args.input}")
        print(f"总页数: {rep['slide_count']}")
        print("=" * 60)
        for s in rep["slides"]:
            flag = "❌ 整页贴图" if s["full_page_image"] else "✅ 真元素"
            mp = s["max_pic_inches"]
            mp_s = f"{mp[0]:.2f}x{mp[1]:.2f}in" if mp else "无"
            print(f"P{s['page']:2d} | 文{s['text_boxes']:2d} 图{s['pictures']:2d} "
                  f"Shape{s['shapes']:2d} | 最大图{mp_s} | {flag}")
        print("=" * 60)
        print("判定:", "FAIL — 存在整页贴图页" if rep["failed"] else "PASS — 无整页贴图")
    sys.exit(1 if rep["failed"] else 0)


if __name__ == "__main__":
    main()
