#!/usr/bin/env python3
"""PPTX 交付前结构自检：检测整页贴图 + 统计可编辑元素。

用法:
    python3 verify_editable.py <file.pptx>            # 检查全部页
    python3 verify_editable.py <file.pptx> --page 6   # 只检查第6页

判定规则:
    - 整页图 = PICTURE 且 宽 > 85% 页面宽 AND 高 > 85% 页面高
    - 存在整页图 → 判不合格，退出码 1（禁止交付）

依赖: python-pptx (pip install python-pptx)
"""
import sys
from pptx import Presentation
from pptx.util import Emu


def check(path, page_only=None):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    fail = 0
    print(f"文件: {path}")
    print(f"页数: {len(prs.slides)} | 页面 {Emu(sw).inches:.2f} x {Emu(sh).inches:.2f} in")
    print("=" * 60)
    for i, slide in enumerate(prs.slides, 1):
        if page_only and i != page_only:
            continue
        shapes = list(slide.shapes)
        pics = [s for s in shapes if s.shape_type == 13]
        txts = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]
        groups = [s for s in shapes if s.shape_type == 6]
        fullpage = [s for s in pics if s.width > sw * 0.85 and s.height > sh * 0.85]
        ok = not fullpage
        if not ok:
            fail += 1
        flag = "OK" if ok else "FAIL 整页贴图"
        print(f"P{i:2d} | 图{len(pics)} 文{len(txts)} 组{len(groups)} | {flag}")
        for s in fullpage:
            print(f"      [全页图] {s.name} {Emu(s.width).inches:.1f}x{Emu(s.height).inches:.1f}in")
    print("=" * 60)
    if fail:
        print(f"FAIL: {fail} 页存在整页贴图，禁止交付（需改为真文本框 + 独立图片对象）")
        return 1
    print("PASS: 无整页贴图，全部为可编辑元素")
    return 0


if __name__ == "__main__":
    args = sys.argv[1:]
    if not args:
        print(__doc__)
        sys.exit(2)
    path = args[0]
    page = None
    if "--page" in args:
        page = int(args[args.index("--page") + 1])
    sys.exit(check(path, page))
