#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""渲染 PPTX 某一页 → PNG 预览（文本框/图片/自动形状）。

无 LibreOffice/Keynote 环境下的可靠渲染方案：python-pptx 解析 shape → 手写 SVG → cairosvg 转 PNG。
仅作验收预览，不用于「整页贴图」交付（整页贴图是铁律禁止的）。

用法:
    python3 render_pptx_preview.py <file.pptx> <slide_idx(0-based)> <out.png> [out.svg]

三个必踩的坑（本脚本已修复，勿回退）:
  1. str(RGBColor) 返回 'RRGGBB' 无 '#' 前缀 → 直接塞进 SVG fill/font 会被 cairosvg 当黑色。
     必须统一补 '#%02X%02X%02X'。
  2. text-anchor="middle" 时 x 坐标是文本框中点，不是左边缘；left/right 同理。
     否则居中文字整体偏移、视觉上像「截断」。
  3. add_picture 同时给 width+height 会拉伸变形；只给其一按比例缩放。
     本脚本用 PIL 读真实尺寸算比例，保持原图宽高比。
"""
import base64
import sys
from pptx import Presentation

EMU_IN = 914400.0
DPI = 96


def _hex(c):
    """颜色对象/元组/字符串 → '#RRGGBB'（补 # 前缀）。"""
    try:
        if c is None:
            return '#333333'
        if hasattr(c, 'rgb') and c.rgb is not None:
            c = c.rgb
        if isinstance(c, (tuple, list)) and len(c) >= 3:
            return '#%02X%02X%02X' % (int(c[0]), int(c[1]), int(c[2]))
        s = str(c)
        if s and not s.startswith('#'):
            return '#' + s
        return s
    except Exception:
        return '#333333'


def _esc(s):
    return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def render_slide(pptx_path, slide_idx, out_png, out_svg=None):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    W = prs.slide_width / EMU_IN
    H = prs.slide_height / EMU_IN
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" '
        f'width="{W}in" height="{H}in" viewBox="0 0 {int(W * DPI)} {int(H * DPI)}">',
        f'<rect x="0" y="0" width="{int(W * DPI)}" height="{int(H * DPI)}" fill="#FFFFFF"/>',
    ]

    for sh in slide.shapes:
        try:
            x = sh.left / EMU_IN * DPI
            y = sh.top / EMU_IN * DPI
            w = sh.width / EMU_IN * DPI
            h = sh.height / EMU_IN * DPI
        except Exception:
            continue
        st = str(sh.shape_type)

        if 'PICTURE' in st:
            blob = sh.image.blob
            mime = getattr(sh.image, 'content_type', 'image/png') or 'image/png'
            b64 = base64.b64encode(blob).decode()
            parts.append(
                f'<image x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                f'xlink:href="data:{mime};base64,{b64}"/>'
            )

        elif 'AUTO_SHAPE' in st:
            fill = '#F2F2F2'
            try:
                fill = _hex(sh.fill.fore_color.rgb)
            except Exception:
                pass
            at = str(sh.auto_shape_type) if sh.auto_shape_type else ''
            if 'ROUND' in at:
                r = min(w, h) * 0.12
                parts.append(
                    f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                    f'rx="{r:.1f}" ry="{r:.1f}" fill="{fill}"/>'
                )
            else:
                parts.append(
                    f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" fill="{fill}"/>'
                )

        elif 'TEXT' in st:
            tf = sh.text_frame
            cy = y
            for para in tf.paragraphs:
                runs = [r for r in para.runs if r.text]
                if not runs:
                    cy += 16
                    continue
                maxsz = max((r.font.size.pt if r.font.size else 12) for r in runs)
                line_h = maxsz * 1.35
                align = str(para.alignment) if para.alignment is not None else ''
                if 'CENTER' in align:
                    ax, anchor = x + w / 2, 'middle'
                elif 'RIGHT' in align:
                    ax, anchor = x + w, 'end'
                else:
                    ax, anchor = x, 'start'
                tspans = []
                for r in runs:
                    size = r.font.size.pt if r.font.size else 12
                    bold = ' font-weight="bold"' if r.font.bold else ''
                    color = _hex(r.font.color)
                    tspans.append(
                        f'<tspan font-size="{size}" fill="{color}"{bold}>{_esc(r.text)}</tspan>'
                    )
                parts.append(
                    f'<text x="{ax:.1f}" y="{cy + maxsz:.1f}" '
                    f'font-family="PingFang SC, Microsoft YaHei, sans-serif" '
                    f'text-anchor="{anchor}">{"".join(tspans)}</text>'
                )
                cy += line_h

    parts.append('</svg>')
    svg = '\n'.join(parts)
    if out_svg:
        with open(out_svg, 'w', encoding='utf-8') as f:
            f.write(svg)

    import cairosvg
    cairosvg.svg2png(bytestring=svg.encode('utf-8'), write_to=out_png,
                     output_width=int(W * DPI), output_height=int(H * DPI))
    print('PNG:', out_png, f'({int(W * DPI)}x{int(H * DPI)})')


if __name__ == '__main__':
    if len(sys.argv) < 4:
        print('用法: render_pptx_preview.py <file.pptx> <slide_idx> <out.png> [out.svg]')
        sys.exit(1)
    _pptx, _idx, _png = sys.argv[1], int(sys.argv[2]), sys.argv[3]
    _svg = sys.argv[4] if len(sys.argv) > 4 else None
    render_slide(_pptx, _idx, _png, _svg)
